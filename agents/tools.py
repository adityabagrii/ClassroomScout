import io
import os
import textwrap
from typing import Any, Dict, List

import pypandoc
from langchain_core.tools import tool
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

from .logging import log_msg
from .text_utils import cache_path_for_drive

ALLOWED_MIMES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.google-apps.document",
}


def paginate_list(call_fn, key: str, page_size: int) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    token = None
    while True:
        resp = call_fn(pageSize=page_size, pageToken=token).execute()
        out.extend(resp.get(key, []))
        token = resp.get("nextPageToken")
        if not token:
            break
    return out


def build_tools(classroom, drive, docs, db, cache_dir: str):
    @tool
    def classroom_list_courses(page_size: int = 100, course_states: str = "ACTIVE") -> str:
        """List courses. Returns JSON string."""
        items = paginate_list(
            lambda pageSize, pageToken=None: classroom.courses().list(
                pageSize=pageSize, pageToken=pageToken, courseStates=course_states
            ),
            "courses",
            page_size,
        )
        return __import__("json").dumps({"courses": items}, ensure_ascii=False)

    @tool
    def classroom_list_announcements(course_id: str, page_size: int = 50) -> str:
        """List announcements for a course. Returns JSON string."""
        items = paginate_list(
            lambda pageSize, pageToken=None: classroom.courses().announcements().list(
                courseId=course_id, pageSize=pageSize, pageToken=pageToken
            ),
            "announcements",
            page_size,
        )
        return __import__("json").dumps({"announcements": items}, ensure_ascii=False)

    @tool
    def classroom_list_coursework(course_id: str, page_size: int = 50) -> str:
        """List courseWork for a course. Returns JSON string."""
        items = paginate_list(
            lambda pageSize, pageToken=None: classroom.courses().courseWork().list(
                courseId=course_id, pageSize=pageSize, pageToken=pageToken
            ),
            "courseWork",
            page_size,
        )
        return __import__("json").dumps({"courseWork": items}, ensure_ascii=False)

    @tool
    def classroom_list_materials(course_id: str, page_size: int = 50) -> str:
        """List courseWorkMaterial for a course. Returns JSON string."""
        items = paginate_list(
            lambda pageSize, pageToken=None: classroom.courses().courseWorkMaterials().list(
                courseId=course_id, pageSize=pageSize, pageToken=pageToken
            ),
            "courseWorkMaterial",
            page_size,
        )
        return __import__("json").dumps({"courseWorkMaterial": items}, ensure_ascii=False)

    @tool
    def db_get_recent_events(course_id: str, since_ts: int) -> str:
        """Fetch recent events since timestamp. Returns JSON string."""
        rows = db.get_recent_events(course_id, since_ts)
        return __import__("json").dumps({"events": rows}, ensure_ascii=False)

    @tool
    def db_get_event(event_id: str) -> str:
        """Get one event by event_id. Returns JSON string."""
        row = db.get_event(event_id)
        return __import__("json").dumps({"event": row}, ensure_ascii=False)

    @tool
    def db_set_checkpoint(key: str, value: str) -> str:
        """Set a checkpoint value. Returns JSON string."""
        db.set_checkpoint(key, value)
        return __import__("json").dumps({"ok": True}, ensure_ascii=False)

    @tool
    def db_get_checkpoint(key: str, default: str = "0") -> str:
        """Get checkpoint value. Returns JSON string."""
        v = db.get_checkpoint(key, default)
        return __import__("json").dumps({"value": v}, ensure_ascii=False)

    @tool
    def drive_get_metadata(file_id: str) -> str:
        """Drive file metadata. Returns JSON string."""
        fields = "id,name,mimeType,modifiedTime,createdTime,size,md5Checksum,webViewLink,webContentLink"
        meta = drive.files().get(fileId=file_id, fields=fields, supportsAllDrives=True).execute()
        return __import__("json").dumps({"meta": meta}, ensure_ascii=False)

    @tool
    def drive_download(file_id: str, out_path: str) -> str:
        """Download a drive file. Returns JSON string."""
        req = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
        fh = io.FileIO(out_path, "wb")
        dl = MediaIoBaseDownload(fh, req)
        done = False
        while not done:
            _, done = dl.next_chunk()
        return __import__("json").dumps({"path": out_path}, ensure_ascii=False)

    @tool
    def drive_export_google_doc(file_id: str, out_path: str, mime: str = "application/pdf") -> str:
        """Export a Google Doc to given mime type. Returns JSON string."""
        req = drive.files().export_media(fileId=file_id, mimeType=mime)
        fh = io.FileIO(out_path, "wb")
        dl = MediaIoBaseDownload(fh, req)
        done = False
        while not done:
            _, done = dl.next_chunk()
        return __import__("json").dumps({"path": out_path}, ensure_ascii=False)

    @tool
    def file_count_and_extract(path: str, mime_type: str) -> str:
        """Count pages/slides and extract text. Returns JSON string."""
        count = 0
        text = ""

        if mime_type == "application/pdf":
            reader = PdfReader(path)
            count = len(reader.pages)
            text = "\n".join([(p.extract_text() or "") for p in reader.pages])

        elif "presentation" in mime_type or mime_type == "application/vnd.ms-powerpoint":
            prs = Presentation(path)
            count = len(prs.slides)
            parts = []
            for s in prs.slides:
                for sh in s.shapes:
                    if hasattr(sh, "text") and sh.text:
                        parts.append(sh.text)
            text = "\n".join(parts)

        elif "wordprocessingml" in mime_type or mime_type == "application/msword":
            doc = Document(path)
            words = sum(len((p.text or "").split()) for p in doc.paragraphs)
            count = max(1, (words + 299) // 300)
            text = "\n".join([p.text for p in doc.paragraphs if p.text])

        return __import__("json").dumps({"count": int(count), "text": text}, ensure_ascii=False)

    @tool
    def docs_create(title: str) -> str:
        """Create a Google Doc. Returns JSON string."""
        doc = docs.documents().create(body={"title": title}).execute()
        return __import__("json").dumps({"documentId": doc["documentId"]}, ensure_ascii=False)

    @tool
    def docs_append(document_id: str, text: str) -> str:
        """Append text to a Google Doc. Returns JSON string."""
        if not text.endswith("\n"):
            text += "\n"
        reqs = [{"insertText": {"location": {"index": 1}, "text": text}}]
        docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        return __import__("json").dumps({"ok": True}, ensure_ascii=False)

    @tool
    def notify_whatsapp(message: str) -> str:
        """Notify user via WhatsApp (stub). Returns JSON string."""
        log_msg(f"notify.whatsapp message={message[:900]}")
        return __import__("json").dumps({"ok": True}, ensure_ascii=False)

    @tool
    def drive_upload(file_path: str, mime_type: str = "application/pdf", name: str = None) -> str:
        """Upload a file to Drive. Returns JSON string."""
        file_metadata = {"name": name or os.path.basename(file_path)}
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        f = drive.files().create(body=file_metadata, media_body=media, fields="id,webViewLink").execute()
        return __import__("json").dumps({"id": f["id"], "webViewLink": f.get("webViewLink")}, ensure_ascii=False)

    def render_markdown_to_pdf(md_path: str, pdf_path: str) -> bool:
        """Render markdown to PDF via pandoc/xelatex."""
        header_tex = "pandoc_wrap_header.tex"
        header = r"""
        % --- Better line breaking for long URLs/strings ---
        \usepackage{xurl}
        \usepackage{microtype}
        \setlength{\emergencystretch}{3em}
        \Urlmuskip=0mu plus 1mu

        % --- Wrap code blocks produced by Pandoc ---
        \usepackage{fvextra}
        \DefineVerbatimEnvironment{Highlighting}{Verbatim}{
        breaklines=true,
        breakanywhere=true,
        commandchars=\\\{\}
        }
        \fvset{breaklines=true,breakanywhere=true}
        """

        with open(header_tex, "w", encoding="utf-8") as f:
            f.write(textwrap.dedent(header).strip() + "\n")

        import sys
        mono = "Menlo" if sys.platform == "darwin" else None

        extra_args = [
            "--pdf-engine=xelatex",
            "-V",
            "geometry:margin=1in",
            "--include-in-header",
            header_tex,
        ]

        if mono:
            extra_args += ["-V", f"monofont={mono}"]
        extra_args += ["-V", "monofontoptions=Scale=0.90"]

        try:
            pypandoc.convert_file(md_path, "pdf", outputfile=pdf_path, extra_args=extra_args)
            return os.path.exists(pdf_path)
        except Exception as e:
            log_msg(f"pdf.render.failed error={e}")
            return False

    def safe_cache_path_for_drive(file_id: str, mime_type: str) -> str:
        return cache_path_for_drive(cache_dir, file_id, mime_type)

    return {
        "classroom_list_courses": classroom_list_courses,
        "classroom_list_announcements": classroom_list_announcements,
        "classroom_list_coursework": classroom_list_coursework,
        "classroom_list_materials": classroom_list_materials,
        "db_get_recent_events": db_get_recent_events,
        "db_get_event": db_get_event,
        "db_set_checkpoint": db_set_checkpoint,
        "db_get_checkpoint": db_get_checkpoint,
        "drive_get_metadata": drive_get_metadata,
        "drive_download": drive_download,
        "drive_export_google_doc": drive_export_google_doc,
        "file_count_and_extract": file_count_and_extract,
        "docs_create": docs_create,
        "docs_append": docs_append,
        "notify_whatsapp": notify_whatsapp,
        "drive_upload": drive_upload,
        "render_markdown_to_pdf": render_markdown_to_pdf,
        "safe_cache_path_for_drive": safe_cache_path_for_drive,
    }, ALLOWED_MIMES
